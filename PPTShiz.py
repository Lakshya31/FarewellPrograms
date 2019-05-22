import pptx
import random

SLD_LAYOUT_TITLE_AND_CONTENT = 1

prs = pptx.Presentation()
title_slide_layout = prs.slide_layouts[0]

List = []
Order = []

for i in range(0,150):
    List.append(i)

random.shuffle(List)

TDS = []

f = open("This is for TDS.txt","r")
for i in range(150):
    _temp = f.readline()
    _temp = _temp[0:-1]
    TDS.append(_temp)



for i in range(150):
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    _temp = List[i]
    if _temp < 50:
        _temp2 = "#"+str(i+1)+"\nOut Comes The TRUTH!"
    elif _temp <100:
        _temp2 = "#"+str(i+1)+"\nThe DARE is ON!"
    else:
        _temp2 = "#"+str(i+1)+"\nSomeone has a SITUATION incoming!"


    title.text = _temp2
    subtitle.text = TDS[_temp]


prs.save('TDS.pptx')